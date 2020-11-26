#!/usr/bin/env python  
# encoding: utf-8  

""" 
@version: v1.0 
@author: xag 
@license: Apache Licence  
@contact: xinganguo@gmail.com 
@site: http://www.xingag.top 
@software: PyCharm 
@file: ppt_utils.py 
@time: 2020-11-14 15:15 
@description：PPT工具类
"""

import math
import os
import random

from PIL import Image
from moviepy.editor import VideoFileClip
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.shapes.picture import Movie
from pptx.util import Inches, Cm
from pptx.util import Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE



def get_video_frame(clip, frame_index):
    """
    获取视频的某一帧图片
    :param clip:
    :param frame_index:
    :return:
    """
    # 帧数目
    frame_count = math.floor(clip.fps * clip.duration)
    # print('视频帧数目:', frame_count)

    # 保证参数输入有效
    if frame_index < 0 or frame_index > frame_count:
        frame_index = 1

    # 视频所有的帧
    frames = clip.iter_frames()
    # clip.get_frame()

    # 定义输出图片路径
    thumbnail_path = "{}/temp/{}.jpg".format(os.path.abspath(os.path.dirname(__file__)), random_str(10))

    # 遍历，找到对应的帧，保存到本地
    for index, frame in enumerate(frames):
        if frame_index == index:
            # 保持帧图片到本地
            im = Image.fromarray(frame)
            im.save(thumbnail_path)
            break

    return thumbnail_path


def get_image_aspect_ratio(image_path):
    """
    获取图片的宽高比
    :param image_path:
    :return:
    """
    img = Image.open(image_path)

    # 图片类型：GIF
    image_format = img.format

    # 图片宽、高
    width, height = img.size

    # 图片宽高比
    aspect_ratio = width / height

    return aspect_ratio


def get_video_aspect_ratio_and_thumbnail_path(video_path, frame_index):
    """
    获取图片的宽、高比
    :param video_path: 视频路径
    :param frame_index 帧索引
    :return:
    """
    clip = VideoFileClip(video_path)

    # 视频的宽度、高度
    width, height = clip.size

    # 获取宽、高比
    aspect_ratio = width / height

    # 获取视频缩略图
    thumbnail_path = get_video_frame(clip, frame_index)

    return aspect_ratio, thumbnail_path


def add_slide(presentation, slide_style_index):
    """
    在PPT文档中，以内置的版式添加幻灯片
    :param presentation:文档对象
    :param slide_style_index:版式索引
    :return:
    """
    # PPT版式样式
    # 内置有11种版式样式
    # 0：Title Slide 标题幻灯片
    # 1：Title and Content  标题和内容
    # 2：Section Header   节标题
    # 3：Two Content   两栏内容
    # 4：Comparison  比较
    # 5：Title Only  仅标题
    # 6：Blank  空白
    # 7：Content with Caption   内容和标题
    # 8：Picture with Caption   图片和标题
    # 9：Title and Vertical Text  标题和竖排内容
    # 10：Vertical Title and Text  竖排标题和文本
    slide_layout = presentation.slide_layouts[slide_style_index]

    # 通过样式Layout，新增一张幻灯片
    slide = presentation.slides.add_slide(slide_layout)

    return slide


def get_slides(presentation):
    """
    获取所有的幻灯片
    :param presentation:
    :return:
    """
    # 所有幻灯片
    slides = presentation.slides

    # 幻灯片数目
    slide_num = len(slides)

    return slides, slide_num


def get_slide(presentation, slide_index):
    """
    根据索引，获取某一个幻灯片
    :param presentation:
    :param slide_index:页面索引，从0开始
    :return:
    """
    slides, slide_num = get_slides(presentation=presentation)

    return slides[slide_index]


def del_slide(presentation, slide_index=0):
    """
    删除某一张幻灯片
    :param presentation:
    :param slide_index: 索引
    :return:
    """
    # 所有幻灯片的列表
    slides = list(presentation.slides._sldIdLst)

    # 根据索引，删除某一张幻灯片
    presentation.slides._sldIdLst.remove(slides[slide_index])


def insert_textbox(slide, left, top, width, height, unit=Inches):
    """
    幻灯片中添加文本框
    :param unit: 单元，默认设置为Inches
    :param slide: 幻灯片对象
    :param left: 左边距
    :param top:  上边距
    :param width: 宽度
    :param height: 高度
    :return:
    """
    # 文本框
    textbox = slide.shapes.add_textbox(left=unit(left),
                                       top=unit(top),
                                       width=unit(width),
                                       height=unit(height))
    # 文本框形状
    tf = textbox.text_frame

    return textbox, tf


def insert_shape(slide, left, top, width, height, autoshape_type_id=MSO_SHAPE.CHEVRON, unit=Inches):
    """
    幻灯片中添加形状
    :param unit: 单位，默认为Inches
    :param autoshape_type_id: 形状类型
    :param slide:幻灯片
    :param left:左边距
    :param top:上边距
    :param width:宽度
    :param height:高度
    :return:
    """
    # 添加一个形状
    # add_shape(self, autoshape_type_id, left, top, width, height)
    # 参数分别为：形状类型、左边距、上边距、宽度、高度
    shape = slide.shapes.add_shape(autoshape_type_id=autoshape_type_id,
                                   left=unit(left),
                                   top=unit(top),
                                   width=unit(width),
                                   height=unit(height))
    return shape


def insert_table(slide, rows, cols, left, top, width, height, unit=Cm):
    """
    幻灯片中插入一个表格
    :param unit: 默认单位为厘米
    :param slide: 幻灯片对象
    :param rows: 行数
    :param cols: 列数
    :param left: 左边距
    :param top: 上边距
    :param width: 宽度
    :param height: 高度
    :return:
    """
    # 插入一个表格
    # <class 'pptx.shapes.graphfrm.GraphicFrame'>
    table = slide.shapes.add_table(rows, cols, unit(left), unit(top), unit(width), unit(height))

    # 返回表格对象
    # <class 'pptx.table.Table'>
    return table.table


def insert_image(slide, pic_path, left, top, width=None, height=None, unit=Inches):
    """
    幻灯片中加入图片(包含静态图片和动态图片)
    :param unit: 单位默认为Inches
    :param pic_path: 文件路径
    :param slide: 幻灯片对象
    :param left: 左边距
    :param top:  上边距
    :param width: 宽度
    :param height: 高度
    :return:
    """
    # 注意：如果width、height都为None时，以图片原始大小展示
    width = unit(width) if width else None
    height = unit(height) if height else None

    # 单位：英尺
    pic_obj = slide.shapes.add_picture(image_file=pic_path,
                                       left=unit(left),
                                       top=unit(top),
                                       width=width,
                                       height=height)
    return pic_obj


def insert_video(slide, video_path, thumbnail_path, left, top, width, height, unit=Cm):
    """
    插入视频到PPT中
    :param slide:幻灯片对象
    :param video_path:视频路径
    :param thumbnail_path:视频封面图
    :param left:左边距
    :param top:上边距
    :param width:宽度
    :param height:高度
    :param unit:单位，默认是厘米
    :return:
    """
    # 注意：视频封面图可以为空，不指定
    video_obj = slide.shapes.add_movie(video_path, unit(left), unit(top), unit(width), unit(height),
                                       poster_frame_image=thumbnail_path)

    return video_obj


def insert_chart(slide, left, top, width, height, data, unit=Inches, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    """
    插入图表
    :param slide: 幻灯片
    :param left: 左边距
    :param top: 上边距
    :param width: 宽度
    :param height: 高度
    :param data: 图表数据
    :param unit: 数据单位，默认为：Inches
    :param chart_type: 图表类型，默认是：柱状图
    :return:
    """
    chart_result = slide.shapes.add_chart(chart_type=chart_type,
                                          x=unit(left), y=unit(top),
                                          cx=unit(width), cy=unit(height),
                                          chart_data=data)
    # 返回图表
    return chart_result.chart


def set_chart_title(chart, title_content, font_name=None, font_size=-1, font_color=None, font_bold=False,
                    font_italic=False):
    """
    为图表添加一个标题
    :param font_italic:
    :param font_bold:
    :param chart:
    :param title_content: 
    :param font_name: 
    :param font_size: 
    :param font_color: 
    :return: 
    """
    # 通过添加一个段落来添加一个标题
    paragraph = chart.chart_title.text_frame.add_paragraph()
    paragraph.text = title_content

    # 设置段落样式
    set_parg_font_style(paragraph, font_name, font_color, font_size, font_bold, font_italic)


def set_widget_bg(widget, bg_rgb_color=None):
    """
    设置【文本框textbox/单元格/形状】的背景颜色
    :param widget:文本框textbox、单元格、形状
    :param bg_rgb_color:背景颜色值
    :return:
    """
    if bg_rgb_color and len(bg_rgb_color) == 3:
        # 1、将形状填充类型设置为纯色
        widget.fill.solid()
        # 2、设置文本框的背景颜色
        widget.fill.fore_color.rgb = RGBColor(bg_rgb_color[0], bg_rgb_color[1], bg_rgb_color[2])


def set_widget_frame(widget, frame_rgb_color=None, frame_width=-1, unit=Cm):
    """
    设置【文本框textbox/单元格/形状】边框样式
    包含：边框颜色、宽度
    :param unit: 默认单位是厘米
    :param widget:文本框textbox/单元格/形状
    :param frame_rgb_color:
    :param frame_width:
    :return:
    """
    if frame_rgb_color and len(frame_rgb_color) == 3:
        widget.line.color.rgb = RGBColor(frame_rgb_color[0], frame_rgb_color[1], frame_rgb_color[2])
    if frame_width > 0:
        widget.line.width = unit(frame_width)


def set_font_style(font, font_name=None, font_color=None, font_size=-1, font_bold=False, font_italic=False):
    """
    设置字体样式
    :param font:
    :param font_name:
    :param font_color:
    :param font_size:
    :param font_bold:
    :param font_italic:
    :return:
    """
    # 字体名称
    if font_name:
        font.name = font_name

    # 字体颜色
    if font_color and len(font_color) == 3:
        font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])

    # 字体大小
    if font_size != -1:
        font.size = Pt(font_size)

    # 是否加粗，默认不加粗
    font.bold = font_bold

    # 是否倾斜，默认不倾斜
    font.italic = font_italic


def set_parg_font_style(paragraph, font_name=None, font_color=None, font_size=-1, font_bold=False, font_italic=False,
                        paragraph_alignment=PP_ALIGN.CENTER):
    """
    设置段落中文本的样式，包含：字体名称、颜色、大小、是否加粗、是否斜体
    :param paragraph_alignment: 段落对齐方式
    :param paragraph:
    :param font_name:
    :param font_color:
    :param font_size:
    :param font_bold:
    :param font_italic:
    :return:
    """

    # 对齐方式
    # 注意：对齐方式是针对段落的
    paragraph.alignment = paragraph_alignment

    # 获取段落中字体对象
    font = paragraph.font

    # 设置字体样式
    set_font_style(font, font_name, font_color, font_size, font_bold, font_italic)

    return font


def set_cell_center(cell):
    """
    设置单元格文字居中显示
    :param cell:
    :return:
    """
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_table_column_width(table, column_index, width, unit=Cm):
    """
    设置表格某一列的宽度
    :param table:
    :param column_index:
    :param width:
    :param unit: 单位默认为厘米
    :return:
    """
    table.columns[column_index].width = unit(width)


def set_table_row_height(table, row_index, height, unit=Cm):
    """
    设置表格某一行的高度
    :param table:
    :param row_index:
    :param height:
    :param unit:
    :return:
    """
    table.rows[row_index].height = unit(height)


def random_str(num):
    """
    生成随机的字符串
    :param num:
    :return:
    """
    result = random.sample('ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789', num)
    return ''.join(result)


def save_ppt_images(presentation, output_path):
    """
     保存ppt中所有图片
    :param presentation:
    :param output_path 保存目录
    :return:
    """

    print('幻灯片数目:', len(presentation.slides))

    # 遍历所有幻灯片
    for index_slide, slide in enumerate(presentation.slides):
        # 遍历所有形状
        for index_shape, shape in enumerate(slide.shapes):
            # 形状包含：文字形状、图片、普通形状等

            # 过滤出图片形状
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 获取图片二进制字符流
                image_data = shape.image.blob

                # image/jpeg、image/png等
                image_type_pre = shape.image.content_type

                # 图片后缀名
                image_suffix = image_type_pre.split('/')[1]

                # 创建image文件夹保存抽出图片
                if not os.path.exists(output_path):
                    os.makedirs(output_path)

                # 图片保存路径
                output_image_path = output_path + random_str(10) + "." + image_suffix

                print(output_image_path)

                # 写入到新的文件中
                with open(output_image_path, 'wb') as file:
                    file.write(image_data)


def get_shape_content(shape, shape_type=0):
    """
    获取存在文本框（has_text_frame为True）的形状Shape的文本内容
    :param shape_type:形状类型
    :param shape:形状
    :return:
    """
    # 文本框：shape.text_frame
    # 文本框中的段落：shape.text_frame.paragraphs
    # 文本框中的文本内容：shape.text_frame.text
    # 文字块Run：shape.text_frame.paragraphs[0].runs
    if shape_type == 0:  # 全部文本
        return shape.text_frame.text
    elif shape_type == 1:  # 按照段落去获取内容
        return [paragraph.text for paragraph in shape.text_frame.paragraphs]
    else:  # 通过文字块，去获取全部内容
        # 形状中所有段落
        paragraphs = shape.text_frame.paragraphs
        # 文字块内容列表
        results = []
        for paragraph in paragraphs:
            runs = paragraph.runs
            for run in runs:
                results.append(run.text)

        return results


def read_ppt_content(presentation):
    """
    读取PPT中所有的内容
    :param presentation:
    :return:
    """
    # 所有内容
    results = []

    # 遍历所有幻灯片，获取文本框中的值
    for slide in presentation.slides:
        for shape in slide.shapes:
            # 判断形状是否包含文本框
            if shape.has_text_frame:
                content = get_shape_content(shape)
                if content:
                    results.append(content)

    return results



