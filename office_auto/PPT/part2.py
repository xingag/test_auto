#!/usr/bin/env python  
# encoding: utf-8  

""" 
@version: v1.0 
@author: xag 
@license: Apache Licence  
@contact: xinganguo@gmail.com 
@site: http://www.xingag.top 
@software: PyCharm 
@file: part2.py 
@time: 2020-11-21 16:27 
@description：Python操作PPT（2）
"""

from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_ANCHOR
from ppt_utils import *


class PptF(object):

    def __init__(self):
        self.presentation = Presentation()

    def start(self):
        # 1.表格管理
        # self.table_manage()

        # 2、图片管理
        # self.image_manage()

        # 3、视频管理
        self.video_manage()

        self.tear_down()

    def tear_down(self):
        self.presentation.save('./output.pptx')

    def table_manage(self):
        """
        表格管理
        :return:
        """
        # 1.创建一个幻灯片 Slide 对象（空白样式)
        slide = add_slide(self.presentation, 6)

        # 2.插入一个表格（3行3列）
        # 参数分别为：幻灯片对象、行数、列数、左边距、上边距、宽度、高度
        table = insert_table(slide, 3, 3, 3, 5, 13.6, 5)

        # 3.重新设置表的宽度、高度
        # 3.1 分别设置第1-3行列宽
        set_table_column_width(table, 0, 5)
        set_table_column_width(table, 1, 5)
        set_table_column_width(table, 2, 5)

        # 3.2 分别设置行高
        set_table_row_height(table, 0, 1.5)
        set_table_row_height(table, 1, 1.2)
        set_table_row_height(table, 2, 1.2)

        # 4.设置表格数据
        datas = [
            ["学员", "姓名", "年龄"],
            ["", "星安果", 23],
            ["", "AirPython", 18]]

        # 遍历设置数据到单元格中
        for row_index in range(len(table.rows)):
            for column_index in range(len(table.columns)):
                # 获取单元格对象
                cell_temp = table.cell(row_index, column_index)

                # 设置数据
                cell_temp.text = str(datas[row_index][column_index])

        # 5、设置第一行表头单元格文字加粗居中显示
        for column_index in range(len(table.columns)):
            # 1、单元格对象
            cell = table.cell(0, column_index)
            # 2、文本控件的段落
            paragraph = cell.text_frame.paragraphs[0]
            # 3、设置段落样式
            set_parg_font_style(paragraph, font_name='微软雅黑', font_size=23, font_color=[255, 0, 0],
                                font_bold=True)

            # 设置单元格背景颜色
            set_widget_bg(cell, [204, 217, 225])

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 6、单元格合并
        # 合并单元格并居中显示
        table.cell(1, 0).merge(table.cell(2, 0))
        table.cell(1,0).text="合并"
        set_cell_center(table.cell(1,0))

    def image_manage(self):
        """
        图片管理
        :return:
        """
        # 插入一张静态图片和一张动态图片
        slide = add_slide(self.presentation, 6)

        # 图片路径
        image_path = './1.jpeg'
        # image_path = './700.gif'

        # 插入本地图片（可能会出现展示不全的情况）
        # insert_image(slide, image_path, 6, 6, unit=Cm)

        # 获取宽、高比
        aspect_ratio = get_image_aspect_ratio(image_path)

        # 等比例插入gif图片到PPT中
        insert_image(slide, image_path, 6, 6, 6, 6 / aspect_ratio, unit=Cm)

    def video_manage(self):
        """
        视频管理
        :return:
        """
        slide = add_slide(self.presentation, 6)

        video_path = './1.mp4'

        # 获取图片宽高比，并保存一个临时的缩略图到本地
        aspect_ratio, thumbnail_path = get_video_aspect_ratio_and_thumbnail_path(video_path, 120)

        # 将视频插入到PPT中
        insert_video(slide, video_path, thumbnail_path, 3, 3, 4, 4 / aspect_ratio)


if __name__ == '__main__':
    pptF = PptF()
    pptF.start()
