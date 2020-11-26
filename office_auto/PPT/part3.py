#!/usr/bin/env python  
# encoding: utf-8  

""" 
@version: v1.0 
@author: xag 
@license: Apache Licence  
@contact: xinganguo@gmail.com 
@site: http://www.xingag.top 
@software: PyCharm 
@file: part3.py 
@time: 2020-11-21 20:58 
@description：Python操作PPT（3）
"""

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_utils import *


class PptF(object):

    def __init__(self):
        self.presentation = Presentation()

    def start(self):
        # 1、预设形状管理
        # self.shape_manage()

        # 2、图表Chart管理
        # self.chart_manage()

        # 3、读取PPT文字内容
        # self.word_manage()

        # 4、保存所有图片到本地
        self.save_images()

        # self.tear_down()

    def tear_down(self):
        self.presentation.save('./output.pptx')

    def shape_manage(self):
        """
        形状管理
        :return:
        """
        slide = add_slide(self.presentation, 6)

        # 1、添加一个圆角矩形
        rectangle = insert_shape(slide, 2, 2, 16, 8, autoshape_type_id=MSO_SHAPE.ROUNDED_RECTANGLE, unit=Cm)

        # 2、设置形状属性
        # 2.1 背景颜色
        set_widget_bg(rectangle, bg_rgb_color=[255, 255, 255])

        # 2.2 边框属性
        set_widget_frame(rectangle, frame_rgb_color=[255, 0, 0], frame_width=0.5)

    def chart_manage(self):
        """
        图表管理
        :return:
        """
        slide = add_slide(self.presentation, 6)

        # 创建一个图表数据对象
        chart_data = ChartData()

        # 数据类别（x轴数据）
        chart_data.categories = [2000, 2005, 2010, 2015, 2020]

        # 每一年各维度的数据（3个纬度）
        # 经济
        chart_data.add_series("经济", [60, 65, 75, 90, 95])

        # 环境
        chart_data.add_series("环境", [95, 88, 84, 70, 54])

        # 文化
        chart_data.add_series("军事", [40, 65, 80, 95, 98])

        # 添加图表
        chart = insert_chart(slide, 4, 5, 20, 9, chart_data, unit=Cm, chart_type=XL_CHART_TYPE.LINE)

        # 设置图表显示属性
        # 显示图例
        chart.has_legend = True

        # 图例是否在绘图区之外显示
        chart.legend.include_in_layout = False

        # 设置图表是否显示平滑
        chart.series[0].smooth = True
        chart.series[1].smooth = True
        chart.series[2].smooth = True

        # 设置图表中文字的样式
        set_font_style(chart.font, font_size=12, font_color=[255, 0, 0])

    def read_ppt_file_table(self):
        """
        读取PPT中的数据
        :return:
        """
        # 打开待读取的ppt
        presentation = Presentation("./raw.pptx")

        for slide in presentation.slides:
            # 遍历素有形状
            # 形状：有内容的形状、无内容的形状
            for shape in slide.shapes:
                # print('当前形状名称:', shape.shape_type)
                # 只取表格中的数据，按照行读取内容
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    # 获取表格行（shape.table.rows）
                    for row in shape.table.rows:
                        # 某一行所有的单元格(row.cells)
                        for cell in row.cells:
                            # 单元格文本框中的内容(cell.text_frame.text)
                            print(cell.text_frame.text)

    def word_manage(self):
        """
        读取PPT中的文本内容
        :return:
        """
        presentation = Presentation("./raw.pptx")

        # 1、普通形状内容的所有文本内容
        contents = read_ppt_content(presentation)
        print(contents)

        # 2、图标中的所有的文本内容
        self.read_ppt_file_table()

    def save_images(self):
        """
        保存PPT中所有图片到本地
        :return:
        """
        presentation = Presentation("./raw.pptx")

        # 保存PPT中所有图片到本地
        save_ppt_images(presentation, './output/')


if __name__ == '__main__':
    pptF = PptF()
    pptF.start()
