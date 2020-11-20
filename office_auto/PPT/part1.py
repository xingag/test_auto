#!/usr/bin/env python  
# encoding: utf-8  

""" 
@version: v1.0 
@author: xag 
@license: Apache Licence  
@contact: xinganguo@gmail.com 
@site: http://www.xingag.top 
@software: PyCharm 
@file: part1.py 
@time: 2020-11-18 23:39 
@description：Python操作PPT（1）
"""

from pptx import Presentation

from ppt_utils import *


class PptF(object):

    def __init__(self):
        self.presentation = Presentation()

    def start(self):
        # 1、PPT幻灯片管理
        # self.ppt_manage()

        # 2、文字和段落
        self.word_manage()

        self.tear_down()

    def ppt_manage(self):
        """
        PPT幻灯片管理
        :return:
        """

        # 内置11种模板
        # 当然，也可以利用占位符（placeholder）自定义母版，插入数据

        # 1.1 新增幻灯片
        slide1 = add_slide(self.presentation, 0)
        slide2 = add_slide(self.presentation, 1)
        slide3 = add_slide(self.presentation, 2)
        slide4 = add_slide(self.presentation, 3)

        # 1.2.1 获取所幻灯片
        slides, slide_num = get_slides(self.presentation)
        print('现有幻灯片:', slides)
        print('幻灯片数目:', slide_num)

        # 1.2.2 获取某一个幻灯片
        slide = get_slide(self.presentation, 1)
        print(slide.shapes)

        # 1.3 根据索引，删除PPT文档中某一张幻灯片
        # 比如：删除第4张幻灯片
        del_slide(self.presentation, 3)

    def word_manage(self):
        """
        文字和段落
        :return:
        """
        # 1、创建一张空白的幻灯片
        slide = add_slide(self.presentation, 6)

        # 2、往幻灯片中插入一个文本框，返回一个文本框对象和一个文本框形状对象
        textbox, tf = insert_textbox(slide, 8, 2, 10, 4, unit=Cm)

        # 2.1 默认的段落
        paragraph_default = tf.paragraphs[0]
        paragraph_default.text = "设置段落默认的内容设置段落默认的内容设置段落默认的内容设置段落默认的内容设置段落默认的内容设置段落默认的内容设置段落默认的内容"

        # 2.2 添加一个新的段落
        paragraph_new = tf.add_paragraph()

        # 2.3 给段落设置内容
        paragraph_new.text = "欢迎关注公众号：AirPython\n每周分享 Python 原创技术干货！"

        # 3、设置段落对齐方式及字体样式
        set_parg_font_style(paragraph_new, font_color=[255, 0, 0])

        # 4、设置文字框的背景颜色
        set_widget_bg(textbox, [0, 255, 0])

        # 5、设置文本框的文字自动对齐
        tf.word_wrap = True


    def tear_down(self):
        self.presentation.save('./output.pptx')


if __name__ == '__main__':
    pptF = PptF()
    pptF.start()

